from pptx import Presentation


def build_ppt_from_outline(ppt_outline: dict, output_path: str):
    """
    ppt_outline = {
      "title": "...",
      "slides": [
        { "title": "...", "bullets": [...] }
      ]
    }
    """
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = ppt_outline["title"]

    for slide_data in ppt_outline.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data.get("title", "")

        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.clear()

        for bullet in slide_data.get("bullets", []):
            tf.add_paragraph().text = bullet

    prs.save(output_path)
